"""删除PPT内的备注和图片广告"""

from pathlib import Path

from pptx import Presentation
from python_pptx_text_replacer import TextReplacer
from tqdm import tqdm

from log import logger


def __fun_替换本文(ppt_file: str) -> None:
    """替换PPT里面的文本."""
    text_need_replace = ["唐峰", "芒果", "T500"]

    replacer = TextReplacer(
        presentation_file_name=ppt_file,
        slides="",
        tables=True,
        charts=True,
        textframes=True,
    )
    for text in text_need_replace:
        replacer.replace_text(replacements=[(text, "小夕")])

    replacer.write_presentation_to_file(presentation_output_file_name=ppt_file)


def __fun_删除所有备注和广告图片(ppt_file: str, ad_pic_name_list: list[str]) -> None:
    """删除PPT里面的备注和广告图片."""
    msg = f"处理PPT:{ppt_file}"
    logger.info(msg=msg)
    prs = Presentation(pptx=ppt_file)

    for slide in prs.slides:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide.notes_slide.notes_text_frame.clear()

        try:
            for shape in slide.shapes:
                if shape.shape_type:
                    if shape.shape_type.name == "PICTURE":
                        for ad_name in ad_pic_name_list:
                            if f'descr="{ad_name}"' in shape.element.xml:
                                slide.shapes._spTree.remove(shape._element)  # noqa: SLF001
                                logger.info(msg="删除广告图片")

                    # 如果是视频
                    # 直接干掉
                    elif shape.shape_type.name == "MEDIA":
                        slide.shapes._spTree.remove(shape._element)  # noqa: SLF001

        except Exception as e:  # noqa: BLE001
            logger.info(e)

    prs.save(file=ppt_file)


def fun_处理所有PPT(material_path: str) -> None:
    """处理所有的PPT"""
    ad_pic_name_list = [f"{x}-({y})" for x in range(1, 100) for y in range(1, 100)]
    ad_pic_name_list.extend(
        [
            f"{start_stem}-({x})"
            for x in range(1, 100)
            for start_stem in ["tm", "mm", "aa"]
        ],
    )

    for in_file in tqdm(
        list(Path(material_path).rglob("*")),
        ncols=100,
        desc="删除PPT备注",
    ):
        if in_file.is_file() and in_file.suffix.lower() in [".ppt", ".pptx"]:
            pic_path = in_file.with_suffix(suffix=".png")
            if pic_path.exists() is False:
                try:
                    __fun_删除所有备注和广告图片(
                        ppt_file=in_file.as_posix(),
                        ad_pic_name_list=ad_pic_name_list,
                    )
                    __fun_替换本文(ppt_file=in_file.as_posix())
                except:  # noqa: E722
                    logger.info(msg="error")
